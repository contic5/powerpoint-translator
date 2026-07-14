from pptx import Presentation
from deep_translator import GoogleTranslator
import time
import pandas as pd
import os

MAX_SHAPES_TO_TRANSLATE=999
TARGET_LANGUAGE="Spanish"
private_mode=True

def translate_presentation(target_file):
    if not ".pptx" in target_file:
        target_file=target_file+".pptx"
    print(f"Translating {target_file}")

    out_file=target_file.replace(".pptx",f"_{TARGET_LANGUAGE}.pptx")
    if(out_file.count("_")==1):
        out_file=out_file.replace("_"," ")
    
    if private_mode:
        prs = Presentation(f'original_private/{target_file}')
    else:
        prs = Presentation(f'original/{target_file}')

    total_shapes=0
    print(f"Total Slides {len(prs.slides)}")
    for slide in prs.slides:
        if total_shapes>=MAX_SHAPES_TO_TRANSLATE:
            break
        
        for shape in slide.shapes:
            #Try to translate shape text
            try:
                translated_text=GoogleTranslator(source='auto', target='es').translate(shape.text)
                print(translated_text)
                shape.text=translated_text
                time.sleep(1)
                total_shapes+=1

                if total_shapes>=MAX_SHAPES_TO_TRANSLATE:
                    break
            #Skip if the shape has no text
            except AttributeError:
                pass

    print(f"Total Shapes {total_shapes}")
    print(f"{target_file} translated to {TARGET_LANGUAGE}")

    if private_mode:
        prs.save(f'output_private/{out_file}')
        os.open(f'output_private/{out_file}')
    else:
        prs.save(f'output/{out_file}')
        os.open(f'output/{out_file}')

def get_language_code(TARGET_LANGUAGE):
    language_iso_df=pd.read_excel("Language_ISO_Codes.xlsx")
    target_language_row=language_iso_df[language_iso_df["Language"]==TARGET_LANGUAGE]
    for index,row in target_language_row.iterrows():
        return row["ISO Code"]
    return ""

def main():
    language_code=get_language_code(TARGET_LANGUAGE)
    print(language_code)

    print("Powerpoint Presentation Translator")
    print("Put PowerPoint files you want to translate in the 'original' folder")
    while True:
        target_file=input("Enter the file name you want to translate (Leave this blank to exit): ")
        target_file=target_file.replace("\n","")
        if target_file=="" or target_file=="EXIT" or len(target_file)<2:
            print("Goodbye")
            break
        translate_presentation(target_file)
        

if __name__=="__main__":
    main()