import tkinter as tk
import customtkinter
from tkinterdnd2 import TkinterDnD, DND_FILES
from tkinter import filedialog

from pptx import Presentation
from deep_translator import GoogleTranslator
import time
import os
import sys
import subprocess

#How many shapes to translate in the Google Translator API
MAX_SHAPES_TO_TRANSLATE=999

#TO DO: Add a way of selecting language

TARGET_LANGUAGE="Spanish"

class App(TkinterDnD.Tk):
    def __init__(self):
        super().__init__()
        self.geometry("450x250")
        self.title("PowerPoint Translator")

        #The label for the title
        self.title_label=customtkinter.CTkLabel(self,text="PowerPoint Translator")
        self.title_label.cget("font").configure(size=24)
        self.title_label.pack()

        #A button that takes files by drag and drop or by click upload
        self.button = customtkinter.CTkButton(
            self,
            text="➕\nUpload File Here\n (Click or Drag and Drop)",
            bg_color="blue",
            fg_color="white",
            command=self.import_file,
            width=200,
            height=50
        )
        self.button.pack(padx=40, pady=20)

        self.button.drop_target_register(DND_FILES)
        self.button.dnd_bind("<<Drop>>", self.drop)

        #Label that tracks translation progress
        self.progress_label=customtkinter.CTkLabel(self,text="Upload a PowerPoint file to translate it.")
        self.progress_label.pack()

        #Powerpoint file
        self.prs=None

        #How many shapes have been translated
        self.shapes_translated=0

        #How many shapes there are across all slides
        self.total_shapes=0

        self.slide_index=0
        self.total_slides=0

        #The target file name
        self.file_name=""

        #The output file name
        self.out_file=""

    #Import from file path on the computer
    def import_file(self):
        file_path = filedialog.askopenfilename(title="Select a file", filetypes=[("PowerPoint Files", "*.pptx"),])
        if file_path:
            print("Selected file:", file_path)
            self.translate_powerpoint(file_path)

    #Import from drag and drop file
    def drop(self, event):
        file_path=event.data
        file_path=file_path.replace("{","")
        file_path=file_path.replace("}","")
        self.button.configure(text=file_path)
        print("Selected file:", file_path)
        self.translate_powerpoint(file_path)

    #Count how many shapes there across all slides
    def get_total_shapes(self):
        total_shapes=0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                total_shapes+=1
        return total_shapes

    def translate_shape(self):
        slide=self.prs.slides[self.slide_index]
        shape=slide.shapes[self.shape_index]

        #Assume the shape was not translated
        shape_was_translated=False

        #If there is text to translate and it was translated correctly, set translated to true.
        try:
            translated_text=GoogleTranslator(source='auto', target='es').translate(shape.text)
            shape.text=translated_text
            shape_was_translated=True
        except AttributeError:
            pass

        #Show progress
        status_text=f'''
        Progress
        Slides: {self.slide_index+1}/{self.total_slides}
        Shapes: {self.shapes_translated}/{self.total_shapes}
        '''
        self.progress_label.configure(text=status_text)

        self.move_to_next_shape(slide,shape_was_translated)
    
    def move_to_next_shape(self,slide,shape_was_translated):
        self.shape_index+=1
        self.shapes_translated+=1

        #If we are out of shapes in the current slide, move to the next slide
        if self.shape_index>=len(slide.shapes):
            #Quietly save presentation in the background after each slide is finished.
            self.save_presentation()
            self.slide_index+=1
            self.shape_index=0

        #If we finished translating the last slide or translated the maximum number of shapes, we are done.
        if self.slide_index>=len(self.prs.slides) or self.shapes_translated>=MAX_SHAPES_TO_TRANSLATE:
            #Final presentation save
            self.after(1000, self.save_presentation,True)
        else:
            #If a shape was translated, wait a second before calling the translator again. Otherwise, call the translator.
            if shape_was_translated:
                self.after(1000, self.translate_shape)
            else:
                self.translate_shape()
            
    #Open the presentation in PowerPoint
    def open_presentation(self):
        if sys.platform == "win32":
            os.startfile(f"{os.getcwd()}/results/{self.out_file}")  # Best for Windows
        elif sys.platform == "darwin":
            subprocess.Popen(["open", f"results/{self.out_file}"])  # Best for macOS
        else:
            subprocess.Popen(["xdg-open", f"results/{self.out_file}"])  # Best for Linux
        
    #Save the PowerPoint Presentation
    def save_presentation(self,opening=False):
        os.makedirs("results", exist_ok=True)
        self.prs.save(f"results/{self.out_file}")
        
        if opening:
            print(f"Total Shapes {self.total_shapes}")
            print(f"{self.file_name} translated to {TARGET_LANGUAGE}")
            self.progress_label.configure(text=f"{self.file_name} translated to {TARGET_LANGUAGE}")
            self.after(3000,self.open_presentation)
        else:
            print(f"Saving {self.file_name} in the background")
            
    #Endpoint for upload and for drag and drop
    def translate_powerpoint(self,file_path):
        self.prs = Presentation(file_path)

        #Get the file name
        self.file_name=file_path.split("/")[-1]        
    
        #Add the language to the output file name. Presentation.pptx would become Presentation Spanish.pptx
        self.out_file=self.file_name.replace(".pptx",f"_{TARGET_LANGUAGE}.pptx")
        if(self.out_file.count("_")==1):
            self.out_file=self.out_file.replace("_"," ")
        print(self.file_name,self.out_file)
        print("Powerpoint Opened")

        #Reset important presentation variables
        self.shapes_translated=0
        self.shape_index=0
        self.total_shapes=self.get_total_shapes()

        self.slide_index=0
        self.total_slides=len(self.prs.slides)

        #Begin translating shapes
        self.translate_shape()

app = App()
app.mainloop()